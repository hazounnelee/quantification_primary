#!/usr/bin/env python3
"""Gate Review PPT 생성기 — 전구체 입자 형상 자동 정량화."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ── 색상 ──────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
ORANGE = RGBColor(0xE8, 0x5D, 0x26)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF2, 0xF5)
DGRAY  = RGBColor(0x2B, 0x2B, 0x2B)
MGRAY  = RGBColor(0x88, 0x88, 0x88)
LBLUE  = RGBColor(0xD6, 0xE4, 0xF7)
LORG   = RGBColor(0xFD, 0xE8, 0xD8)
LGREEN = RGBColor(0xE8, 0xF5, 0xE9)
GREEN  = RGBColor(0x1B, 0x5E, 0x20)
LRED   = RGBColor(0xFD, 0xEB, 0xEB)
RED    = RGBColor(0xC6, 0x28, 0x28)

SW = Inches(13.33)
SH = Inches(7.5)
FN = "맑은 고딕"

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH


# ── 헬퍼 ─────────────────────────────────────────────────────────────────────

def _slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return s


def _rect(s, l, t, w, h, fill=None, line=None, lw=Pt(1)):
    sh = s.shapes.add_shape(1, l, t, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh


def _text(s, text, l, t, w, h, sz=Pt(13), bold=False,
          color=DGRAY, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = sz
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FN
    return tb


def _box(s, lines, l, t, w, h, fill=NAVY, fg=WHITE,
         sz=Pt(13), bold=False, align=PP_ALIGN.CENTER,
         border=None, bw=Pt(1.5)):
    sh = _rect(s, l, t, w, h, fill=fill, line=border, lw=bw)
    tf = sh.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines if isinstance(lines, list) else [lines]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = line
        r.font.size = sz
        r.font.bold = bold
        r.font.color.rgb = fg
        r.font.name = FN
    return sh


def _header(s, title, section=""):
    _rect(s, 0, 0, SW, Inches(0.72), fill=NAVY)
    _rect(s, 0, Inches(0.72), SW, Inches(0.045), fill=ORANGE)
    _text(s, title, Inches(0.4), Inches(0.10), Inches(10), Inches(0.60),
          sz=Pt(20), bold=True, color=WHITE)
    if section:
        _text(s, section, Inches(10.2), Inches(0.19), Inches(2.9), Inches(0.38),
              sz=Pt(10), color=RGBColor(0xCC, 0xD9, 0xF5), align=PP_ALIGN.RIGHT)
    _rect(s, 0, Inches(7.33), SW, Inches(0.02), fill=NAVY)


def _pgnum(s, n, total=20):
    _text(s, f"{n}  /  {total}", Inches(12.1), Inches(7.30),
          Inches(1.0), Inches(0.18), sz=Pt(9), color=MGRAY, align=PP_ALIGN.RIGHT)


def _bullets(s, items, l, t, w, h, sz=Pt(13), color=DGRAY, gap=Pt(5)):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = gap
        depth, text = item if isinstance(item, tuple) else (0, item)
        prefix = "▪  " if depth == 0 else "    –  "
        r = p.add_run()
        r.text = prefix + text
        r.font.size = sz
        r.font.color.rgb = color
        r.font.name = FN
    return tb


def _arrow_v(s, cx, y1, y2, color=MGRAY, w=Pt(1.5)):
    conn = s.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, cx, y1, cx, y2)
    conn.line.color.rgb = color
    conn.line.width = w
    spPr = conn._element.spPr
    # Reuse the existing a:ln element (created by conn.line.color.rgb above)
    # to avoid a duplicate a:ln that would cause tailEnd to be ignored.
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'arrow')
    te.set('w', 'med')
    te.set('len', 'med')
    return conn


def _table(s, rows, l, t, w, h, hfill=NAVY, hfg=WHITE, alt=LGRAY, fsz=Pt(12)):
    tbl = s.shapes.add_table(len(rows), len(rows[0]), l, t, w, h).table
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = fsz
            p.font.name = FN
            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = hfg
                cell.fill.solid()
                cell.fill.fore_color.rgb = hfill
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt
    return tbl



# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 1: 표지
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
bg = s.background.fill
bg.solid()
bg.fore_color.rgb = NAVY

_rect(s, 0, 0, Inches(0.35), SH, fill=ORANGE)
_rect(s, Inches(0.35), 0, Inches(0.06), SH, fill=RGBColor(0x2A, 0x4E, 0x8A))

_text(s, "전구체 입자 형상 자동 정량화",
      Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.3),
      sz=Pt(38), bold=True, color=WHITE)
_text(s, "SEM 이미지 기반 1차 · 2차 입자 품질 판정 시스템",
      Inches(0.9), Inches(3.1), Inches(11), Inches(0.75),
      sz=Pt(20), color=RGBColor(0xCC, 0xD9, 0xF5))
_rect(s, Inches(0.9), Inches(3.9), Inches(5.5), Inches(0.05), fill=ORANGE)

_text(s, "[부서명]  |  [발표자]                                        [날짜]",
      Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5),
      sz=Pt(13), color=RGBColor(0xAA, 0xBB, 0xDD))
_pgnum(s, 1)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 2: Agenda
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "Agenda")
_pgnum(s, 2)

sections = [
    ("01", "문제 정의",   "전구체 품질 판정의 필요성 및 정량화 목표",        "p. 3–4"),
    ("02", "데이터셋",    "소입경 / 대입경 SEM 이미지 구성",               "p. 5"),
    ("03", "방법론",      "파이프라인 · SAM2 · LSD · 구현 환경",           "p. 6–11"),
    ("04", "실험 결과",   "정성적 / 정량적 결과 · 오류 케이스",             "p. 12–16"),
    ("05", "향후 계획",   "미분/깨짐 정량화 · 판상 · SAM2 개선",           "p. 17–18"),
]

y0 = Inches(0.9)
row_h = Inches(1.02)
for i, (num, title, desc, pg) in enumerate(sections):
    y = y0 + i * row_h
    is_odd = (i % 2 == 0)
    bg_fill = LBLUE if is_odd else LGRAY
    _rect(s, Inches(0.35), y, SW - Inches(0.7), row_h - Inches(0.06),
          fill=bg_fill, line=None)
    _box(s, [num], Inches(0.35), y, Inches(0.65), row_h - Inches(0.06),
         fill=NAVY, fg=WHITE, sz=Pt(22), bold=True)
    _text(s, title, Inches(1.15), y + Inches(0.12), Inches(3.2), row_h - Inches(0.15),
          sz=Pt(16), bold=True, color=NAVY)
    _text(s, desc, Inches(4.5), y + Inches(0.22), Inches(7.0), row_h - Inches(0.25),
          sz=Pt(12), color=DGRAY)
    _text(s, pg, Inches(12.2), y + Inches(0.22), Inches(0.85), row_h - Inches(0.25),
          sz=Pt(11), color=MGRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 3: 문제 정의 — 배경 및 필요성
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "문제 정의 — 배경 및 필요성", "01 문제 정의")
_pgnum(s, 3)

_text(s, "전구체(Precursor) 입자 형상은 배터리 용량·수명에 직결되는 핵심 품질 지표입니다.",
      Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.45),
      sz=Pt(13), color=MGRAY)

cards = [
    (LRED,   RED,   "현행 방식의 한계",
     ["수작업 SEM 이미지 분석 → 분석자 편차 존재",
      "수백 장 이미지 처리 시 수 시간 소요",
      "정량적 수치 없이 정성적 판단에 의존",
      "Lot 간 비교·추세 분석 불가"]),
    (LBLUE,  NAVY,  "자동화의 필요성",
     ["일관된 기준으로 대량 이미지 자동 처리",
      "형상 지표(구형도, 두께, 종횡비) 수치화",
      "시계열 품질 변화 추적 가능",
      "공정 조건 최적화 피드백 루프 구축"]),
    (LGREEN, GREEN, "기대 효과",
     ["분석 시간 ○○% 단축",
      "Lot 합격·불합격 기준 객관화",
      "반응 조건 변화에 따른 입자 형상 정량 비교",
      "연구 데이터 DB화"]),
]

for i, (bg, fg, title, items) in enumerate(cards):
    xl = Inches(0.35) + i * Inches(4.28)
    _rect(s, xl, Inches(1.4), Inches(4.1), Inches(5.7), fill=bg, line=fg, lw=Pt(1.5))
    _box(s, [title], xl, Inches(1.4), Inches(4.1), Inches(0.55),
         fill=fg, fg=WHITE, sz=Pt(13), bold=True)
    _bullets(s, items, xl + Inches(0.12), Inches(2.05), Inches(3.85), Inches(5.0),
             sz=Pt(12), color=DGRAY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 4: 문제 정의 — 정량화 목표
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "문제 정의 — 정량화 목표", "01 문제 정의")
_pgnum(s, 4)

_text(s, "본 연구는 SEM 이미지로부터 아래 두 가지 입자 유형에 대한 형상 지표를 자동 추출합니다.",
      Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.45), sz=Pt(13), color=MGRAY)

# 1차입자 박스
_rect(s, Inches(0.5), Inches(1.45), Inches(5.7), Inches(5.5),
      fill=LBLUE, line=NAVY, lw=Pt(2))
_box(s, ["1차 입자 (대입경 · 소입경 구형 전구체)"],
     Inches(0.5), Inches(1.45), Inches(5.7), Inches(0.65),
     fill=NAVY, fg=WHITE, sz=Pt(14), bold=True)
_bullets(s, [
    "구형도 (Sphericity)",
    (1, "Wadell 2D: 4π × area / perimeter²"),
    (1, "1.0 = 완전 구형, < 1.0 = 불규칙"),
    "등가 직경 (Equivalent Diameter)",
    (1, "2√(area / π)  →  µm 변환"),
    "미분 비율 (Fine Particle Ratio)",
    (1, "면적 임계값 미만 fragment 비율 (%)"),
], Inches(0.7), Inches(2.2), Inches(5.3), Inches(4.5),
   sz=Pt(13), color=DGRAY, gap=Pt(7))

# 2차입자 박스
_rect(s, Inches(7.0), Inches(1.45), Inches(5.7), Inches(5.5),
      fill=LORG, line=ORANGE, lw=Pt(2))
_box(s, ["2차 입자 (개별 결정 — 침상 검출)"],
     Inches(7.0), Inches(1.45), Inches(5.7), Inches(0.65),
     fill=ORANGE, fg=WHITE, sz=Pt(14), bold=True)
_bullets(s, [
    "두께 (Thickness)",
    (1, "LSD 수직 프로파일 7샘플 중앙값"),
    (1, "µm 변환"),
    "장축 (Long Axis)",
    (1, "minAreaRect 장변 길이  →  µm"),
    "종횡비 (Aspect Ratio)",
    (1, "두께 / 장축  (0 ~ 1)"),
    "▷ 판상(Plate) 정량화는 향후 계획",
], Inches(7.2), Inches(2.2), Inches(5.3), Inches(4.5),
   sz=Pt(13), color=DGRAY, gap=Pt(7))

# 화살표 & 라벨
_text(s, "측정\n기준", Inches(6.1), Inches(3.6), Inches(1.1), Inches(0.9),
      sz=Pt(12), color=MGRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 5: 데이터셋
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "데이터셋", "02 데이터셋")
_pgnum(s, 5)

_text(s, "[데이터 경로] 의 전구체 SEM 이미지 데이터셋",
      Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4), sz=Pt(13), color=MGRAY)

_table(s, [
    ["구분",       "입자 유형",   "배율",    "해상도",         "이미지 수",  "비고"],
    ["소입경",    "1차 입자",   "20,000×", "2048 × 1636",   "[N]장",     "전처리: 하단 크롭 (기본 50px)"],
    ["대입경",    "1차 입자",   "3,000×",  "2048 × 1636",   "[N]장",     ""],
    ["2차 입자",  "결정 침상",  "20,000×", "2048 × 1636",   "[N]장",     "고배율 상세 분석"],
    ["2차 입자",  "결정 침상",  "50,000×", "2048 × 1636",   "[N]장",     "소입경 침상 상세"],
],
Inches(0.4), Inches(1.38), Inches(12.5), Inches(2.5), fsz=Pt(12))

_text(s, "이미지 전처리 공통 파이프라인",
      Inches(0.4), Inches(4.05), Inches(5), Inches(0.4),
      sz=Pt(13), bold=True, color=NAVY)
prep_steps = [
    ("원본 SEM 이미지\n(가변 해상도)", LGRAY, DGRAY),
    ("리사이즈  (--preprocess_width W)\n기본 W=1024  →  1024 × 818", LBLUE, NAVY),
    ("하단 크롭  (스케일바 제거)\n기본 round(W·100/2048)=50 px", LBLUE, NAVY),
    ("출력 (기본)\n1024 × 768", LORG, ORANGE),
]
bw, bh = Inches(2.4), Inches(0.72)
xstart = Inches(0.4)
for i, (label, bg, fg) in enumerate(prep_steps):
    xl = xstart + i * (bw + Inches(0.55))
    _box(s, [label], xl, Inches(4.55), bw, bh, fill=bg, fg=fg, sz=Pt(11))
    if i < len(prep_steps) - 1:
        _text(s, "→", xl + bw + Inches(0.06), Inches(4.72), Inches(0.42), Inches(0.4),
              sz=Pt(20), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

_text(s, "※ 스케일 기준 (default --preprocess_width 1024):  20k ×  →  74 px = 1 µm     |     50k ×  →  185 px = 1 µm",
      Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.4),
      sz=Pt(12), color=MGRAY)

_text(s, "※ 이미지 파일명 컨벤션: [IMG_ID]_[seq].[ext]  →  동일 IMG_ID 그룹으로 배치 집계",
      Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.4),
      sz=Pt(12), color=MGRAY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 6: 전체 파이프라인 개요
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "전체 파이프라인 개요", "03 방법론")
_pgnum(s, 6)

# 왼쪽: 1차입자 (secondary_measure.py)
_box(s, ["1차 입자 파이프라인", "(secondary_measure.py)"],
     Inches(0.3), Inches(0.82), Inches(5.8), Inches(0.62),
     fill=NAVY, fg=WHITE, sz=Pt(13), bold=True)

steps_1 = [
    ("SEM 이미지 입력", ""),
    ("전처리 (리사이즈 · 크롭)", ""),
    ("ROI 추출", ""),
    ("SAM2 타일 추론", "Shi-Tomasi 포인트 프롬프트"),
    ("마스크 IoU 중복 제거", "bbox(0.85) + pixel(0.60)"),
    ("등가직경 · 구형도 계산", "2√(A/π)  |  4πA/P²"),
    ("CSV / JSON / 히스토그램 저장", ""),
]

y_s = Inches(1.52)
bw_p = Inches(5.5)
bh_p = Inches(0.55)
gap_p = Inches(0.18)
cx1 = Inches(0.4) + bw_p / 2

for i, (label, sub) in enumerate(steps_1):
    fc = NAVY if i == 0 else (LBLUE if i % 2 == 0 else WHITE)
    fg_c = WHITE if fc == NAVY else NAVY
    lines = [label, sub] if sub else [label]
    bh = bh_p + (Inches(0.18) if sub else 0)
    _box(s, lines, cx1 - bw_p / 2, y_s, bw_p, bh,
         fill=fc, fg=fg_c, sz=Pt(11), border=NAVY, bw=Pt(1))
    if i < len(steps_1) - 1:
        _arrow_v(s, cx1, y_s + bh, y_s + bh + gap_p)
    y_s += bh + gap_p

# 오른쪽: 2차입자 (primary_measure.py)
_box(s, ["2차 입자 파이프라인", "(primary_measure.py)"],
     Inches(7.2), Inches(0.82), Inches(5.8), Inches(0.62),
     fill=ORANGE, fg=WHITE, sz=Pt(13), bold=True)

steps_2 = [
    ("SEM 이미지 입력", ""),
    ("전처리 (리사이즈 · 크롭)", ""),
    ("ROI 추출 / 구 자동 검출", ""),
    ("LSD 또는 SAM2 분기", "particle_mode에 따라"),
    ("[LSD]  침상 검출 · 융합", "CLAHE → 이진화 → 선분 필터"),
    ("두께 측정", "수직 프로파일 7샘플 중앙값"),
    ("CSV / JSON / 오버레이 저장", ""),
]

y_s = Inches(1.52)
cx2 = Inches(7.2) + bw_p / 2

for i, (label, sub) in enumerate(steps_2):
    fc = ORANGE if i == 0 else (LORG if i % 2 == 0 else WHITE)
    fg_c = WHITE if fc == ORANGE else DGRAY
    lines = [label, sub] if sub else [label]
    bh = bh_p + (Inches(0.18) if sub else 0)
    _box(s, lines, cx2 - bw_p / 2, y_s, bw_p, bh,
         fill=fc, fg=fg_c, sz=Pt(11), border=ORANGE, bw=Pt(1))
    if i < len(steps_2) - 1:
        _arrow_v(s, cx2, y_s + bh, y_s + bh + gap_p, color=ORANGE)
    y_s += bh + gap_p

# 중앙 구분선
_rect(s, Inches(6.5), Inches(0.85), Inches(0.32), Inches(6.5),
      fill=LGRAY, line=None)
_text(s, "공통\n전처리", Inches(6.5), Inches(3.6), Inches(0.32), Inches(0.9),
      sz=Pt(9), color=MGRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 7: SAM2 소개
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "SAM2 (Segment Anything Model 2) 소개", "03 방법론")
_pgnum(s, 7)

_text(s, "Meta AI, 2024 — 이미지/비디오 범용 Segmentation 모델",
      Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4), sz=Pt(13), color=MGRAY)

# 왼쪽: 구조 설명
_box(s, ["아키텍처 개요"],
     Inches(0.35), Inches(1.38), Inches(5.8), Inches(0.5),
     fill=NAVY, fg=WHITE, sz=Pt(13), bold=True)
_bullets(s, [
    "Image Encoder: Hiera (Hierarchical Vision Transformer)",
    (1, "이미지를 다중 해상도 feature map 으로 인코딩"),
    "Prompt Encoder",
    (1, "포인트 · 박스 · 마스크 프롬프트를 임베딩으로 변환"),
    "Mask Decoder",
    (1, "feature + 프롬프트 임베딩 → 3개 후보 마스크 출력"),
    (1, "각 마스크에 confidence score 부여"),
    "Memory Module (비디오용)",
    (1, "본 프로젝트에서는 이미지 모드만 사용"),
], Inches(0.5), Inches(1.98), Inches(5.5), Inches(4.8),
   sz=Pt(12), color=DGRAY, gap=Pt(6))

# 오른쪽: 적용 방식
_box(s, ["본 프로젝트 적용 방식"],
     Inches(7.0), Inches(1.38), Inches(5.8), Inches(0.5),
     fill=ORANGE, fg=WHITE, sz=Pt(13), bold=True)
_bullets(s, [
    "타일 기반 추론 (512 × 512, stride 256)",
    (1, "대형 이미지 전체를 타일로 분할 처리"),
    "Shi-Tomasi 포인트 샘플링",
    (1, "코너 특징점을 SAM2 포인트 프롬프트로 사용"),
    (1, "타일당 최대 80포인트, 배치 크기 32"),
    "IoU 기반 마스크 중복 제거",
    (1, "1단계: bbox IoU ≥ 0.85 → 즉시 제거"),
    (1, "2단계: pixel IoU ≥ 0.60 → 최종 제거"),
    "모델: SAM2.1 Hiera-Base+",
], Inches(7.15), Inches(1.98), Inches(5.5), Inches(4.8),
   sz=Pt(12), color=DGRAY, gap=Pt(6))

_rect(s, Inches(6.5), Inches(1.38), Inches(0.15), Inches(5.6), fill=LGRAY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 8: LSD vs Hough Transform
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "침상 검출: LSD vs Hough Transform", "03 방법론")
_pgnum(s, 8)

_text(s, "2차 입자 침상(Acicular) 검출에 LSD를 선택한 이유",
      Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4), sz=Pt(13), color=MGRAY)

_table(s, [
    ["비교 항목",          "Hough Transform",                        "LSD (Line Segment Detector)",           "선택"],
    ["알고리즘",           "파라미터 공간 투표 방식",                   "그라디언트 기반 지지 영역(support region)",  "LSD ✓"],
    ["파라미터 의존성",    "threshold · rho · theta 수동 튜닝 필수",   "자동 파라미터 — 튜닝 불필요",               "LSD ✓"],
    ["검출 결과",          "무한 직선 (ρ, θ) — 끝점 없음",             "끝점 좌표 포함 선분 반환",                   "LSD ✓"],
    ["짧은 선분",          "단편적 선분 검출 어려움",                    "최소 길이(기본 20px) 이상 선분 직접 검출",   "LSD ✓"],
    ["연산 속도",          "고해상도 이미지에서 느림",                   "선형 시간 복잡도 O(n)",                     "LSD ✓"],
    ["선분 융합",          "별도 후처리 필요",                          "Union-Find 기반 자체 융합 구현",             "LSD ✓"],
    ["OpenCV 지원",        "cv2.HoughLinesP",                         "cv2.createLineSegmentDetector(0)",          ""],
], Inches(0.35), Inches(1.38), Inches(12.6), Inches(4.1), fsz=Pt(11))

_text(s, "▷  LSD 중복 제거: 중심 거리 < 12 px  AND  방향각 차이 < 25°  →  짧은 선분 제거 (긴 선분 우선 유지)",
      Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.45), sz=Pt(12), color=NAVY)
_text(s, "▷  두께 측정: 각 선분의 수직 방향으로 7개 샘플 → 중앙값(Median)으로 이상치 제거",
      Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.45), sz=Pt(12), color=NAVY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 9: 전처리 · 스케일 캘리브레이션
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "전처리 · 스케일 캘리브레이션", "03 방법론")
_pgnum(s, 9)

# 캘리브레이션 테이블
_box(s, ["스케일 캘리브레이션"],
     Inches(0.35), Inches(0.9), Inches(5.8), Inches(0.5),
     fill=NAVY, fg=WHITE, sz=Pt(13), bold=True)
_table(s, [
    ["배율",      "기준 픽셀",   "기준 거리",  "µm / pixel",   "적용 대상"],
    ["20,000×",   "147 px",     "1 µm",       "0.00680 µm/px", "소입경 1차 · 2차 입자"],
    ["50,000×",   "371 px",     "1 µm",       "0.00269 µm/px", "소입경 2차 입자 상세"],
    ["3,000×",    "[N] px",     "[N] µm",     "직접 입력",      "대입경 1차 입자"],
],
Inches(0.35), Inches(1.48), Inches(5.8), Inches(2.0), fsz=Pt(12))

_text(s, "변환 공식:   length_µm  =  length_px  ×  (scale_µm / scale_px)",
      Inches(0.4), Inches(3.6), Inches(5.8), Inches(0.45), sz=Pt(12), color=NAVY)

# 오른쪽: CLAHE 설명
_box(s, ["대비 향상 — CLAHE"],
     Inches(7.0), Inches(0.9), Inches(5.8), Inches(0.5),
     fill=ORANGE, fg=WHITE, sz=Pt(13), bold=True)
_bullets(s, [
    "Contrast Limited Adaptive Histogram Equalization",
    (1, "타일 단위 히스토그램 평활화 → 국소 대비 향상"),
    (1, "clip limit 으로 노이즈 증폭 억제"),
    "적용 시점: LSD 전처리 단계에서 Gray 이미지에 적용",
    "이진화 방법 (선택 가능)",
    (1, "--lsd_adaptive_thresh OFF (기본) → Otsu 전역 이진화"),
    (1, "--lsd_adaptive_thresh ON       → Adaptive Gaussian 이진화"),
], Inches(7.15), Inches(1.5), Inches(5.5), Inches(3.2),
   sz=Pt(12), color=DGRAY, gap=Pt(6))

# 하단: 전처리 흐름
_box(s, ["이미지 전처리 흐름 (LSD 분기)"],
     Inches(0.35), Inches(4.2), Inches(12.6), Inches(0.48),
     fill=RGBColor(0x37, 0x47, 0x6F), fg=WHITE, sz=Pt(12), bold=True)

flow = ["Gray 변환", "CLAHE", "이진화\n(Otsu/Adaptive +\n자동 반전)", "LSD 선분 검출",
        "중복 제거\n(dist<12px, Δang<25°)", "두께 측정\n(7샘플 중앙값)", "길이 필터\n(--min_length, 선택: --fuse)"]
bw2, bh2 = Inches(1.55), Inches(0.65)
xs = Inches(0.35)
for i, step in enumerate(flow):
    xl = xs + i * (bw2 + Inches(0.16))
    fc = LBLUE if i % 2 == 0 else WHITE
    _box(s, [step], xl, Inches(4.78), bw2, bh2,
         fill=fc, fg=NAVY, sz=Pt(10), border=NAVY, bw=Pt(1))
    if i < len(flow) - 1:
        _text(s, "→", xl + bw2 + Inches(0.01), Inches(4.9),
              Inches(0.14), Inches(0.4),
              sz=Pt(14), bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 10: 구현 환경
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "구현 환경", "03 방법론")
_pgnum(s, 10)

# SW 환경
_box(s, ["소프트웨어 환경"],
     Inches(0.35), Inches(0.9), Inches(5.9), Inches(0.5),
     fill=NAVY, fg=WHITE, sz=Pt(13), bold=True)
_table(s, [
    ["항목",          "내용"],
    ["Language",      "Python 3.11  (conda env: measure)"],
    ["SAM2",          "Ultralytics SAM2  —  sam2.1_hiera_base_plus"],
    ["Computer Vision","OpenCV 4.x  (LSD, CLAHE, morphology)"],
    ["수치 연산",      "NumPy, SciPy"],
    ["시각화",         "Matplotlib (histogram, pipeline PDF)"],
    ["설정 관리",      "PyYAML  —  configs/presets.yaml · paths.yaml"],
    ["실행 환경",      "macOS / Linux  |  CPU 또는 CUDA 선택 가능"],
],
Inches(0.35), Inches(1.48), Inches(5.9), Inches(3.6), fsz=Pt(12))

# HW 환경
_box(s, ["하드웨어 / 실행 방법"],
     Inches(7.0), Inches(0.9), Inches(5.8), Inches(0.5),
     fill=ORANGE, fg=WHITE, sz=Pt(13), bold=True)
_bullets(s, [
    "하드웨어",
    (1, "GPU: [사용 GPU 모델]  (SAM2 추론 가속)"),
    (1, "RAM: [GB]  |  VRAM: [GB]"),
    "1차 입자 분석 실행",
    (1, "python primary_measure.py --input [경로]"),
    (1, "  --preset acicular_20k"),
    "2차 입자 분석 실행",
    (1, "python secondary_measure.py --input [경로]"),
    (1, "  --preset secondary_20k"),
    "출력물",
    (1, "03_overlay_roi.png  (2× 업스케일 · 라벨 포함)"),
    (1, "objects.csv  |  summary.json  |  histogram.png"),
], Inches(7.15), Inches(1.5), Inches(5.5), Inches(5.2),
   sz=Pt(12), color=DGRAY, gap=Pt(5))

# 배포 구조
_box(s, ["프로젝트 구조 (주요 모듈)"],
     Inches(0.35), Inches(5.2), Inches(5.9), Inches(0.45),
     fill=RGBColor(0x37, 0x47, 0x6F), fg=WHITE, sz=Pt(12), bold=True)
modules = [
    ("primary_measure.py", "1차 입자 CLI 진입점"),
    ("secondary_measure.py", "2차 입자 CLI 진입점"),
    ("services/sam2_service.py", "SAM2 기반 공통 서비스"),
    ("services/primary_particle.py", "침상 검출 서비스"),
    ("utils/lsd.py", "LSD · 두께 측정"),
    ("utils/image.py", "전처리 · 시각화 유틸"),
]
for i, (mod, desc) in enumerate(modules):
    row_bg = LGRAY if i % 2 == 0 else WHITE
    y_row = Inches(5.73) + i * Inches(0.27)
    _rect(s, Inches(0.35), y_row, Inches(5.9), Inches(0.27), fill=row_bg)
    _text(s, mod, Inches(0.45), y_row + Inches(0.04), Inches(2.6), Inches(0.22),
          sz=Pt(10), color=NAVY)
    _text(s, desc, Inches(3.1), y_row + Inches(0.04), Inches(3.1), Inches(0.22),
          sz=Pt(10), color=DGRAY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 11: 2차입자(SAM2) 분류 로직
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "1차 입자 분류 로직 (SAM2 기반)", "03 방법론")
_pgnum(s, 11)

_text(s, "SAM2가 출력한 각 마스크를 면적 임계값 기준으로 particle / fragment 로 분류합니다.",
      Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4), sz=Pt(13), color=MGRAY)

# 분류 흐름도
flow_cls = [
    ("SAM2 출력 마스크", "binary mask per object"),
    ("최소 마스크 면적\n필터 (min_valid_area)", "너무 작은 노이즈 제거"),
    ("컨투어 추출\n(RETR_EXTERNAL)", "가장 큰 contour 선택"),
    ("ROI 경계 마진\n체크 (bbox_edge_margin)", "경계 걸친 객체 제외"),
    ("면적 임계값 비교\n(particle_area_threshold)", "기본값: 1500 px²"),
]
bw3, bh3 = Inches(3.8), Inches(0.65)
cx_m = Inches(3.5)
y_f = Inches(1.42)
for i, (label, sub) in enumerate(flow_cls):
    fc = NAVY if i == 0 else LBLUE
    fg_c = WHITE if fc == NAVY else NAVY
    _box(s, [label, sub], cx_m - bw3 / 2, y_f, bw3, bh3 + Inches(0.12),
         fill=fc, fg=fg_c, sz=Pt(11), border=NAVY, bw=Pt(1))
    if i < len(flow_cls) - 1:
        _arrow_v(s, cx_m, y_f + bh3 + Inches(0.12), y_f + bh3 + Inches(0.32))
    y_f += bh3 + Inches(0.44)

# 분기 결과
_rect(s, cx_m - bw3 / 2, y_f, bw3, Inches(0.5),
      fill=RGBColor(0x37, 0x47, 0x6F), line=None)
_text(s, "area ≥ threshold ?", cx_m - bw3 / 2, y_f, bw3, Inches(0.5),
      sz=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

y_out = y_f + Inches(0.65)
_box(s, ["Particle ✓", "구형도 · 등가직경 계산"],
     cx_m - bw3 / 2 - Inches(1.8), y_out, Inches(3.2), Inches(0.75),
     fill=LGREEN, fg=GREEN, sz=Pt(12), bold=False, border=GREEN, bw=Pt(1.5))
_box(s, ["Fragment ✗", "개수만 집계 (미분 지표)"],
     cx_m + bw3 / 2 - Inches(1.4), y_out, Inches(3.2), Inches(0.75),
     fill=LRED, fg=RED, sz=Pt(12), bold=False, border=RED, bw=Pt(1.5))

_text(s, "YES", cx_m - Inches(3.2), y_f + Inches(0.52),
      Inches(1.0), Inches(0.3), sz=Pt(11), color=GREEN, align=PP_ALIGN.CENTER)
_text(s, "NO", cx_m + Inches(2.1), y_f + Inches(0.52),
      Inches(1.0), Inches(0.3), sz=Pt(11), color=RED, align=PP_ALIGN.CENTER)

# 오른쪽: 측정값 목록
_box(s, ["측정 출력값 (Particle만)"],
     Inches(8.0), Inches(0.9), Inches(5.0), Inches(0.5),
     fill=NAVY, fg=WHITE, sz=Pt(13), bold=True)
_bullets(s, [
    "float_eqDiameterUm    등가 직경 (µm)",
    "float_sphericity          구형도 (0 ~ 1)",
    "float_bboxWidthUm / HeightUm    bbox 크기",
    "int_maskArea                 마스크 면적 (px²)",
    "float_confidence            SAM2 신뢰도",
    "str_category                particle / fragment",
],
Inches(8.15), Inches(1.5), Inches(4.8), Inches(4.5),
sz=Pt(12), color=DGRAY, gap=Pt(7))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 12: 실험 결과 — 1차 입자 정성적
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "실험 결과 — 1차 입자 정성적 결과", "04 실험 결과")
_pgnum(s, 12)

_text(s, "SAM2 세그멘테이션 결과 · 등가원 오버레이 시각화 (2× 업스케일, 라벨 표시)",
      Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4), sz=Pt(13), color=MGRAY)

for i, label in enumerate(["소입경 (20k×)", "대입경 (3k×)"]):
    xl = Inches(0.35) + i * Inches(6.4)
    _rect(s, xl, Inches(1.38), Inches(6.15), Inches(5.2),
          fill=LGRAY, line=NAVY, lw=Pt(1))
    _text(s, f"[ {label} 오버레이 이미지 삽입 ]",
          xl + Inches(0.5), Inches(3.2), Inches(5.0), Inches(0.6),
          sz=Pt(13), color=MGRAY, align=PP_ALIGN.CENTER)
    _box(s, [label], xl, Inches(1.38), Inches(6.15), Inches(0.48),
         fill=NAVY, fg=WHITE, sz=Pt(12), bold=True)

_bullets(s, [
    "초록 마스크: Particle  |  주황 마스크: Fragment",
    "각 particle에 라벨 표시:  d = 등가직경(µm),  S = 구형도",
    "초록 원: 등가직경 기준 원 (--eq_diameter 활성화 시)",
], Inches(0.4), Inches(6.72), Inches(12.5), Inches(0.65),
   sz=Pt(12), color=DGRAY, gap=Pt(2))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 13: 실험 결과 — 2차 입자 정성적
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "실험 결과 — 2차 입자 (침상) 정성적 결과", "04 실험 결과")
_pgnum(s, 13)

_text(s, "LSD 기반 침상 검출 오버레이 · 두께 라벨 시각화 (2× 업스케일)",
      Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4), sz=Pt(13), color=MGRAY)

for i, label in enumerate(["원본 ROI", "이진화 결과", "LSD 검출 결과", "두께 라벨 오버레이"]):
    xl = Inches(0.2) + i * Inches(3.25)
    _rect(s, xl, Inches(1.38), Inches(3.0), Inches(3.6),
          fill=LGRAY, line=NAVY, lw=Pt(1))
    _text(s, f"[ {label} ]",
          xl + Inches(0.15), Inches(2.9), Inches(2.7), Inches(0.5),
          sz=Pt(11), color=MGRAY, align=PP_ALIGN.CENTER)
    _box(s, [label], xl, Inches(1.38), Inches(3.0), Inches(0.42),
         fill=ORANGE, fg=WHITE, sz=Pt(11), bold=True)

_bullets(s, [
    "빨강/파랑 선분: 검출된 침상 선분  |  각 선분에 두께(µm) 라벨 표시",

    "ROI 밀도(roi_density): 이진화 이미지 흰 픽셀 비율 — 침상 밀집도 지표",
], Inches(0.4), Inches(5.1), Inches(12.5), Inches(0.65),
   sz=Pt(12), color=DGRAY, gap=Pt(3))

# LSD 단계별 이미지
_box(s, ["LSD 단계별 디버그 이미지"],
     Inches(0.35), Inches(5.88), Inches(12.6), Inches(0.45),
     fill=RGBColor(0x37, 0x47, 0x6F), fg=WHITE, sz=Pt(12), bold=True)
steps_lsd = ["01 Gray", "02 CLAHE", "03 Binary", "04 All LSD",
             "05 Filtered", "06 Fused", "07 Thickness"]
bw_l = Inches(1.68)
for i, st in enumerate(steps_lsd):
    xl = Inches(0.35) + i * (bw_l + Inches(0.04))
    _box(s, [st], xl, Inches(6.4), bw_l, Inches(0.85),
         fill=LGRAY if i % 2 == 0 else WHITE,
         fg=NAVY, sz=Pt(10), border=NAVY, bw=Pt(0.75))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 14: 실험 결과 — 정량적 통계
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "실험 결과 — 정량적 통계", "04 실험 결과")
_pgnum(s, 14)

_box(s, ["1차 입자 — 구형도 · 등가직경 통계"],
     Inches(0.35), Inches(0.9), Inches(12.6), Inches(0.48),
     fill=NAVY, fg=WHITE, sz=Pt(13), bold=True)
_table(s, [
    ["구분",     "n",   "Mean",   "Median",  "Std",    "Min",    "Max"],
    ["소입경 구형도",  "[N]",  "[0.XX]",  "[0.XX]",  "[0.XX]",  "[0.XX]",  "[0.XX]"],
    ["대입경 구형도",  "[N]",  "[0.XX]",  "[0.XX]",  "[0.XX]",  "[0.XX]",  "[0.XX]"],
    ["소입경 등가직경 (µm)", "[N]", "[X.XX]", "[X.XX]", "[X.XX]", "[X.XX]", "[X.XX]"],
    ["대입경 등가직경 (µm)", "[N]", "[X.XX]", "[X.XX]", "[X.XX]", "[X.XX]", "[X.XX]"],
    ["소입경 미분 비율 (%)", "[N]", "[X.X]",  "[X.X]",  "[X.X]",  "[X.X]",  "[X.X]"],
],
Inches(0.35), Inches(1.45), Inches(12.6), Inches(2.3), fsz=Pt(12))

_box(s, ["2차 입자 — 두께 · 종횡비 통계 (침상)"],
     Inches(0.35), Inches(3.9), Inches(12.6), Inches(0.48),
     fill=ORANGE, fg=WHITE, sz=Pt(13), bold=True)
_table(s, [
    ["구분",           "n",   "Mean",   "Median",  "Std",    "Min",    "Max"],
    ["소입경 두께 (µm)",  "[N]",  "[X.XX]",  "[X.XX]",  "[X.XX]",  "[X.XX]",  "[X.XX]"],
    ["소입경 장축 (µm)",  "[N]",  "[X.XX]",  "[X.XX]",  "[X.XX]",  "[X.XX]",  "[X.XX]"],
    ["소입경 종횡비",     "[N]",  "[0.XX]",  "[0.XX]",  "[0.XX]",  "[0.XX]",  "[0.XX]"],
],
Inches(0.35), Inches(4.45), Inches(12.6), Inches(1.7), fsz=Pt(12))

_text(s, "※ 히스토그램은 각 출력 폴더의 thickness_histogram.png / size_histogram.png 참조",
      Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.4), sz=Pt(11), color=MGRAY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 15: 파라미터 분석
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "파라미터 분석 — LSD 핵심 옵션", "04 실험 결과")
_pgnum(s, 15)

_text(s, "주요 LSD 파라미터 변경에 따른 검출 결과 비교",
      Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4), sz=Pt(13), color=MGRAY)

params = [
    ("--lsd_adaptive_thresh","OFF (기본값)", "ON",
     "조명 불균일 이미지에서 Adaptive Gaussian 이진화 → 검출률 향상"),
]
for i, (opt, def_val, alt_val, note) in enumerate(params):
    y_p = Inches(1.45) + i * Inches(2.6)
    _box(s, [opt], Inches(0.35), y_p, Inches(3.5), Inches(0.48),
         fill=NAVY, fg=WHITE, sz=Pt(13), bold=True)
    # 두 케이스 이미지 박스
    for j, val in enumerate([def_val, alt_val]):
        xl = Inches(0.35) + j * Inches(4.85)
        _rect(s, xl + Inches(3.65), y_p, Inches(4.5), Inches(2.0),
              fill=LGRAY, line=NAVY, lw=Pt(1))
        _text(s, f"[ {val} 결과 이미지 ]",
              xl + Inches(3.65) + Inches(0.5), y_p + Inches(0.85), Inches(3.5), Inches(0.5),
              sz=Pt(11), color=MGRAY, align=PP_ALIGN.CENTER)
        _box(s, [val], xl + Inches(3.65), y_p, Inches(4.5), Inches(0.4),
             fill=(LBLUE if j == 0 else LORG), fg=(NAVY if j == 0 else ORANGE),
             sz=Pt(11), bold=True, border=(NAVY if j == 0 else ORANGE))
    _text(s, f"▷  {note}",
          Inches(0.35), y_p + Inches(2.1), Inches(12.6), Inches(0.38),
          sz=Pt(12), color=NAVY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 16: 오류 케이스
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "오류 케이스 — SAM2 세그멘테이션 실패", "04 실험 결과")
_pgnum(s, 16)

_text(s, "SAM2 출력 마스크에서 관찰된 주요 실패 유형 및 원인 분석",
      Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4), sz=Pt(13), color=MGRAY)

cases = [
    ("과분할 (Over-segmentation)",
     "단일 입자가 여러 마스크로 분할 검출",
     "포인트 밀도 과다 또는 경계 대비 부족",
     "IoU 임계값 조정 / 포인트 간격 확대"),
    ("과통합 (Under-segmentation)",
     "인접 입자들이 하나의 마스크로 병합",
     "입자 간 간격 협소, 대비 낮음",
     "타일 크기 축소 / morphology 후처리"),
    ("경계 마스크 제외",
     "ROI 경계에 걸친 입자 측정 불가",
     "bbox_edge_margin 정책에 의한 의도적 제외",
     "마진 값 조정 또는 경계 입자 별도 처리"),
]

for i, (title, symptom, cause, remedy) in enumerate(cases):
    xl = Inches(0.35) + i * Inches(4.3)
    _rect(s, xl, Inches(1.38), Inches(4.1), Inches(5.7),
          fill=LRED, line=RED, lw=Pt(1.5))
    _box(s, [title], xl, Inches(1.38), Inches(4.1), Inches(0.5),
         fill=RED, fg=WHITE, sz=Pt(12), bold=True)
    _rect(s, xl + Inches(0.1), Inches(1.98), Inches(3.9), Inches(1.8),
          fill=LGRAY, line=None)
    _text(s, "[ 실패 사례 이미지 ]",
          xl + Inches(0.6), Inches(2.55), Inches(3.0), Inches(0.5),
          sz=Pt(11), color=MGRAY, align=PP_ALIGN.CENTER)
    _text(s, "증상", xl + Inches(0.12), Inches(3.9), Inches(0.55), Inches(0.3),
          sz=Pt(10), bold=True, color=RED)
    _text(s, symptom, xl + Inches(0.7), Inches(3.9), Inches(3.3), Inches(0.45),
          sz=Pt(11), color=DGRAY)
    _text(s, "원인", xl + Inches(0.12), Inches(4.45), Inches(0.55), Inches(0.3),
          sz=Pt(10), bold=True, color=RED)
    _text(s, cause, xl + Inches(0.7), Inches(4.45), Inches(3.3), Inches(0.55),
          sz=Pt(11), color=DGRAY)
    _text(s, "대안", xl + Inches(0.12), Inches(5.1), Inches(0.55), Inches(0.3),
          sz=Pt(10), bold=True, color=GREEN)
    _text(s, remedy, xl + Inches(0.7), Inches(5.1), Inches(3.3), Inches(0.65),
          sz=Pt(11), color=GREEN)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 17: 향후 계획 — 1차입자
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "향후 계획 — 1차 입자 (구형 전구체)", "05 향후 계획")
_pgnum(s, 17)

_box(s, ["현재"],
     Inches(0.35), Inches(0.9), Inches(3.0), Inches(0.5),
     fill=NAVY, fg=WHITE, sz=Pt(14), bold=True)
_bullets(s, [
    "SAM2 기반 세그멘테이션",
    "등가직경 · 구형도 자동 산출",
    "미분 비율 (fragment ratio) 집계",
    "소입경 / 대입경 별도 파이프라인",
], Inches(0.35), Inches(1.5), Inches(3.0), Inches(2.0), sz=Pt(12), color=DGRAY)

_text(s, "→", Inches(3.55), Inches(2.2), Inches(0.6), Inches(0.6),
      sz=Pt(28), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

plan_items = [
    ("미분/깨짐 세분화 정량화", LORG, ORANGE,
     ["현재 면적 임계값으로만 구분 → 형상 지표(원형도, 장축비) 기반 세분화",
      "깨짐(Crack) 패턴 감지: 컨투어 오목(Concavity) 분석",
      "미분 입자 크기 분포 별도 히스토그램 출력"]),
    ("반응 조건별 구형도 추세 분석", LBLUE, NAVY,
     ["반응 시간 / 온도 / pH별 구형도 변화 자동 집계",
      "Lot 간 비교 대시보드 (CSV → 시각화 자동화)",
      "이상 Lot 조기 탐지 알림"]),
]

for i, (title, bg, fg, items) in enumerate(plan_items):
    xl = Inches(4.35) + i * Inches(4.3)
    _rect(s, xl, Inches(0.9), Inches(4.1), Inches(6.2),
          fill=bg, line=fg, lw=Pt(2))
    _box(s, [title], xl, Inches(0.9), Inches(4.1), Inches(0.55),
         fill=fg, fg=WHITE, sz=Pt(12), bold=True)
    _bullets(s, items, xl + Inches(0.15), Inches(1.55),
             Inches(3.85), Inches(5.5), sz=Pt(12), color=DGRAY, gap=Pt(8))

_rect(s, Inches(0.35), Inches(3.6), Inches(3.0), Inches(0.02), fill=MGRAY)
_text(s, "장기 목표",
      Inches(0.35), Inches(3.7), Inches(3.0), Inches(0.35),
      sz=Pt(11), bold=True, color=MGRAY)
_bullets(s, [
    "Lot 품질 DB 연계 → 공정 파라미터 최적화 피드백",
    "실시간 인라인 측정 시스템 연동 검토",
], Inches(0.35), Inches(4.1), Inches(3.0), Inches(2.5), sz=Pt(11), color=MGRAY, gap=Pt(6))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 18: 향후 계획 — 2차입자
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "향후 계획 — 2차 입자 (개별 결정)", "05 향후 계획")
_pgnum(s, 18)

_box(s, ["현재"],
     Inches(0.35), Inches(0.9), Inches(3.0), Inches(0.5),
     fill=ORANGE, fg=WHITE, sz=Pt(14), bold=True)
_bullets(s, [
    "LSD 기반 침상(Acicular) 검출",
    "두께 · 장축 · 종횡비 측정",
    "선분 융합(Union-Find) 구현",
    "판상(Plate) 미측정",
], Inches(0.35), Inches(1.5), Inches(3.0), Inches(2.0), sz=Pt(12), color=DGRAY)

_text(s, "→", Inches(3.55), Inches(2.2), Inches(0.6), Inches(0.6),
      sz=Pt(28), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

plan2 = [
    ("판상(Plate) 입자 정량화", LORG, ORANGE,
     ["SAM2 + minAreaRect → 판상 두께 · 넓이 측정",
      "판상/침상 자동 분류 기준 정립",
      "혼합 시료(침상+판상) 동시 분석"]),
    ("SAM2 세그멘테이션 오류 개선", LBLUE, NAVY,
     ["Convex Hull 기반 마스크 보정",
      "  → 오목 결함(over-seg) 교정",
      "외접원(Minimum Enclosing Circle) 폴백",
      "  → 마스크 추출 실패 시 대안 크기 추정",
      "추가 후처리: morphology open/close 강화"]),
    ("LSD 파라미터 자동 최적화", LGREEN, GREEN,
     ["이미지 밝기·대비에 따른 threshold 자동 조정",
      "Adaptive 이진화 블록 크기 최적화",
      "배율별 최소 선분 길이 자동 캘리브레이션"]),
]

for i, (title, bg, fg, items) in enumerate(plan2):
    xl = Inches(4.35) + i * Inches(3.0)
    _rect(s, xl, Inches(0.9), Inches(2.85), Inches(6.2),
          fill=bg, line=fg, lw=Pt(2))
    _box(s, [title], xl, Inches(0.9), Inches(2.85), Inches(0.62),
         fill=fg, fg=WHITE, sz=Pt(11), bold=True)
    _bullets(s, items, xl + Inches(0.12), Inches(1.62),
             Inches(2.65), Inches(5.4), sz=Pt(11), color=DGRAY, gap=Pt(7))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 19: Summary
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
_header(s, "Summary", "")
_pgnum(s, 19)

_text(s, "전구체 입자 형상 자동 정량화 시스템 핵심 성과",
      Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.4), sz=Pt(14), color=MGRAY)

achievements = [
    ("01", "통합 파이프라인 구축",
     "1차 입자(SAM2) · 2차 입자(LSD+SAM2) 파이프라인을\n단일 코드베이스로 통합"),
    ("02", "핵심 정량화 지표",
     "1차: 구형도 · 등가직경 · 미분비율\n2차: 두께 · 장축 · 종횡비"),
    ("03", "LSD 선택 및 개선",
     "Hough 대비 파라미터 불필요, 끝점 포함\nUnion-Find 선분 융합으로 정확도 향상"),
    ("04", "확장성",
     "배율별 preset · paths.yaml 설정 파일\nCLI 인자로 모든 파라미터 제어 가능"),
]

for i, (num, title, body) in enumerate(achievements):
    row = i // 2
    col = i % 2
    xl = Inches(0.35) + col * Inches(6.4)
    yt = Inches(1.42) + row * Inches(2.65)
    _rect(s, xl, yt, Inches(6.2), Inches(2.45), fill=LGRAY, line=NAVY, lw=Pt(1))
    _box(s, [num], xl, yt, Inches(0.62), Inches(2.45),
         fill=NAVY, fg=WHITE, sz=Pt(22), bold=True)
    _text(s, title, xl + Inches(0.75), yt + Inches(0.2), Inches(5.3), Inches(0.5),
          sz=Pt(15), bold=True, color=NAVY)
    _text(s, body, xl + Inches(0.75), yt + Inches(0.75), Inches(5.3), Inches(1.6),
          sz=Pt(12), color=DGRAY)

_rect(s, Inches(0.35), Inches(6.88), SW - Inches(0.7), Inches(0.45), fill=NAVY)
_text(s, "GitHub:  hazounnelee / measure",
      Inches(0.55), Inches(6.9), Inches(12), Inches(0.38),
      sz=Pt(12), color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 20: Appendix — 반응 시간별 입도 변화
# ════════════════════════════════════════════════════════════════════════════
s = _slide()
bg2 = s.background.fill
bg2.solid()
bg2.fore_color.rgb = LGRAY
_rect(s, 0, 0, SW, Inches(0.72), fill=RGBColor(0x37, 0x47, 0x6F))
_rect(s, 0, Inches(0.72), SW, Inches(0.045), fill=ORANGE)
_text(s, "Appendix  —  반응 시간에 따른 입도(등가직경) 증가",
      Inches(0.4), Inches(0.10), Inches(12.5), Inches(0.60),
      sz=Pt(18), bold=True, color=WHITE)
_pgnum(s, 20)

_text(s, "2차 입자(1차 전구체 구형체) 합성 시간에 따른 입자 크기 변화 측정 결과",
      Inches(0.4), Inches(0.88), Inches(12.5), Inches(0.4), sz=Pt(13), color=MGRAY)

# 그래프 영역
_rect(s, Inches(0.35), Inches(1.38), Inches(7.8), Inches(5.7),
      fill=WHITE, line=RGBColor(0x37, 0x47, 0x6F), lw=Pt(1.5))
_text(s, "[ 반응 시간 vs 등가직경 그래프 삽입\n  X축: 반응 시간 (h),  Y축: 등가직경 (µm) ]",
      Inches(1.2), Inches(3.5), Inches(6.0), Inches(1.0),
      sz=Pt(13), color=MGRAY, align=PP_ALIGN.CENTER)

# 오른쪽: 수치 테이블
_box(s, ["시간별 입도 통계"],
     Inches(8.35), Inches(1.38), Inches(4.6), Inches(0.5),
     fill=RGBColor(0x37, 0x47, 0x6F), fg=WHITE, sz=Pt(12), bold=True)
_table(s, [
    ["반응 시간 (h)", "n",   "Mean (µm)", "Std"],
    ["[T1]",         "[N]", "[X.XX]",    "[X.XX]"],
    ["[T2]",         "[N]", "[X.XX]",    "[X.XX]"],
    ["[T3]",         "[N]", "[X.XX]",    "[X.XX]"],
    ["[T4]",         "[N]", "[X.XX]",    "[X.XX]"],
    ["[T5]",         "[N]", "[X.XX]",    "[X.XX]"],
],
Inches(8.35), Inches(1.95), Inches(4.6), Inches(2.5), fsz=Pt(12))

_bullets(s, [
    "반응 시간 증가에 따라 등가직경 단조 증가 관측",
    "T[N] 이후 성장 둔화 → 포화 구간 확인",
    "Lot 간 동일 시간에서의 입도 편차 분석 가능",
], Inches(8.35), Inches(4.6), Inches(4.6), Inches(2.4),
   sz=Pt(12), color=DGRAY, gap=Pt(7))


# ── 저장 ──────────────────────────────────────────────────────────────────────
out = "gate_review_precursor.pptx"
prs.save(out)
print(f"저장 완료: {out}  ({prs.slides.__len__()} 슬라이드)")
